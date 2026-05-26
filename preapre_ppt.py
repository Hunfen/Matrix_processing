#!/Users/hunfen/Documents/GitHub/STM_DataProcessing/.venv/bin/python
"""Generate PowerPoint from already-processed images folder.

Usage: python generate_ppt.py <output_folder>

The output_folder must contain an 'images' subdirectory with Z_mtrx_forward_up PNGs.
The PPT will be saved as <output_folder_name>.pptx in the same folder.
"""

import logging
import re
import sys
from datetime import datetime
from pathlib import Path

import h5py
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

logger = logging.getLogger(__name__)


def natural_key(s: str) -> list:
    """Split the string by numeric parts and convert them to integers for natural sorting."""
    return [
        int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", s)
    ]


def get_z_images_sorted(img_dir: Path) -> list[Path]:
    """从images目录中收集Z_mtrx_forward_up通道的PNG，按自然顺序返回路径列表。"""
    return sorted(
        img_dir.glob("*Z_mtrx_forward_up.png"),
        key=lambda p: natural_key(p.stem),
    )


def parse_sts_filename(filename: str) -> tuple[str, int] | None:
    """从 sts 文件名中提取 image 基础名和组索引，如 'image1_group2.png' -> ('image1', 2)。"""
    m = re.match(r"^(.+)_group(\d+)\.png$", filename)
    if m:
        img_name = m.group(1)
        group_idx = int(m.group(2))
        return img_name, group_idx
    return None


# def generate_ppt(output_folder: Path) -> None:
#     """创建包含首页和Z_mtrx图像页的PPT，保存到output_folder下。

#     Parameters
#     ----------
#     output_folder : Path
#         已生成images子目录的输出目录。
#     """
#     prs = Presentation()

#     # ---------- 首页：使用 output_folder 的父级三层路径和创建时间 ----------
#     slide = prs.slides.add_slide(prs.slide_layouts[6])

#     parts = output_folder.resolve().parts
#     # 取最后三层作为显示名称（类似 test.py 的逻辑）
#     name_folder = "\\".join(parts[-3:]) if len(parts) >= 3 else "\\".join(parts)
#     time_folder = str(datetime.fromtimestamp(output_folder.resolve().stat().st_ctime))
#     time_folder = time_folder.split(".")[0]

#     text_box1 = slide.shapes.add_textbox(Cm(4), Cm(5.5), Inches(4), Inches(1))
#     tf1 = text_box1.text_frame
#     p1 = tf1.add_paragraph()
#     p1.text = name_folder
#     p1.font.size = Pt(44)

#     line = slide.shapes.add_shape(1, Cm(0), Cm(9.5), Cm(25.4), Cm(1.3))
#     line.fill.solid()
#     line.fill.fore_color.rgb = RGBColor(0, 0, 255)
#     line.line.fill.background()

#     text_box2 = slide.shapes.add_textbox(Cm(4), Cm(13), Inches(4), Inches(1))
#     tf2 = text_box2.text_frame
#     p2 = tf2.add_paragraph()
#     p2.text = time_folder
#     p2.font.size = Pt(18)

#     # ---------- 内容页：Z_mtrx图片 ----------
#     img_dir = output_folder / "images"
#     if not img_dir.is_dir():
#         logger.error("Images directory not found: %s", img_dir)
#         sys.exit(1)

#     z_images = get_z_images_sorted(img_dir)
#     if z_images:
#         pic_width = Cm(4.5)
#         pic_height = Cm(4.5)
#         left_margin = Cm(0.5)
#         gap_x = (Cm(25.4) - 2 * left_margin - 5 * pic_width) / 4
#         row1_y = Cm(2)
#         row2_y = Cm(9)
#         label_height = Cm(0.5)
#         label_offset = Cm(0.6)

#         page_count = 0
#         for start in range(0, len(z_images), 10):
#             chunk = z_images[start : start + 10]
#             slide = prs.slides.add_slide(prs.slide_layouts[6])
#             for i, img_path in enumerate(chunk):
#                 col = i % 5
#                 row = i // 5
#                 left = left_margin + col * (pic_width + gap_x)
#                 top = row1_y if row == 0 else row2_y

#                 # 插入图片
#                 slide.shapes.add_picture(
#                     str(img_path), left, top, width=pic_width, height=pic_height
#                 )

#                 # 提取短标签（-- 之后的数字部分）
#                 stem = img_path.stem
#                 label = stem.split("--")[-1].replace(".Z_mtrx_forward_up", "")

#                 # 图片正上方添加标签文本框
#                 txBox = slide.shapes.add_textbox(
#                     left, top - label_offset, pic_width, label_height
#                 )
#                 tf = txBox.text_frame
#                 tf.word_wrap = True
#                 p = tf.paragraphs[0]
#                 p.text = label
#                 p.font.size = Pt(7)
#                 p.alignment = PP_ALIGN.CENTER

#                 # ---------- 新增：图片下方参数信息 ----------
#                 if stem.endswith(".Z_mtrx_forward_up"):
#                     orig_name = stem[: -len(".Z_mtrx_forward_up")]
#                     h5_path = output_folder / (orig_name + ".Z_mtrx.h5")
#                     text_str = ""
#                     if h5_path.exists():
#                         try:
#                             with h5py.File(h5_path, "r") as hf:
#                                 params = hf["parameters"]
#                                 bias = round(
#                                     float(
#                                         params["GapVoltageControl.Voltage"].attrs[
#                                             "value"
#                                         ]
#                                     ),
#                                     3,
#                                 )
#                                 sp_pa = round(
#                                     float(params["Regulator.Setpoint_1"].attrs["value"])
#                                     * 1e12,
#                                     2,
#                                 )
#                                 # 直接从图像属性读取真实尺寸 (m)，并转换为 nm
#                                 img_grp = hf["images"]
#                                 first_dir = next(iter(img_grp))
#                                 w_m = float(img_grp[first_dir].attrs["width"])
#                                 h_m = float(img_grp[first_dir].attrs["height"])
#                                 real_w_nm = w_m * 1e9
#                                 real_h_nm = h_m * 1e9
#                                 text_str = (
#                                     f"{bias} V, {sp_pa} pA\n"
#                                     f"{real_w_nm:.1f} × {real_h_nm:.1f} nm²"
#                                 )
#                         except Exception as e:
#                             logger.warning(
#                                 "Failed to read params from %s: %s", h5_path, e
#                             )
#                     else:
#                         logger.warning("HDF5 not found: %s", h5_path)
#                     if text_str:
#                         txBox_bottom = slide.shapes.add_textbox(
#                             left,
#                             top + pic_height + Cm(0.05),
#                             pic_width,
#                             Cm(0.6),
#                         )
#                         tf_bottom = txBox_bottom.text_frame
#                         tf_bottom.word_wrap = True
#                         p_bottom = tf_bottom.paragraphs[0]
#                         p_bottom.text = text_str
#                         p_bottom.font.size = Pt(6)
#                         p_bottom.alignment = PP_ALIGN.CENTER

#             page_count += 1
#         logger.info("Added %d content pages with Z_mtrx images", page_count)
#     else:
#         logger.warning("No Z_mtrx_forward_up images found in %s", img_dir)
#     # ---------- 新增 STS 与 mask 叠加页 ----------
#     sts_dir = output_folder / "sts"
#     mask_dir = output_folder / "mask"
#     if sts_dir.is_dir() and mask_dir.is_dir():
#         sts_files = sorted(sts_dir.glob("*.png"), key=lambda p: natural_key(p.stem))
#         added_pages = 0
#         for sts_path in sts_files:
#             info = parse_sts_filename(sts_path.name)
#             if info is None:
#                 continue
#             img_name, group_idx = info
#             # 对应的 image 和 mask 路径
#             img_path = img_dir / f"{img_name}_forward_up.png"
#             mask_path = mask_dir / f"{img_name}_group{group_idx}.png"
#             if not img_path.exists() or not mask_path.exists():
#                 logger.warning("Missing image/mask for %s", sts_path.name)
#                 continue

#             slide = prs.slides.add_slide(prs.slide_layouts[6])

#             # 左侧叠加区域尺寸
#             left_img = Cm(1)
#             top_img = Cm(3)
#             img_size = Cm(12)  # 宽高相等，因为原图是方形

#             # 插入 image
#             slide.shapes.add_picture(
#                 str(img_path), left_img, top_img, width=img_size, height=img_size
#             )
#             # 插入 mask 在完全相同位置（透明 PNG 覆盖）
#             slide.shapes.add_picture(
#                 str(mask_path), left_img, top_img, width=img_size, height=img_size
#             )

#             # 右侧 STS 图片
#             sts_left = Cm(14)
#             sts_top = Cm(3)
#             sts_width = Cm(9)
#             sts_height = Cm(9)
#             slide.shapes.add_picture(
#                 str(sts_path), sts_left, sts_top, width=sts_width, height=sts_height
#             )

#             # 可选的组标签
#             txBox = slide.shapes.add_textbox(left_img, Cm(1.5), img_size, Cm(1))
#             tf = txBox.text_frame
#             p = tf.paragraphs[0]
#             p.text = f"{img_name}  Group {group_idx}"
#             p.font.size = Pt(14)
#             p.alignment = PP_ALIGN.LEFT

#             added_pages += 1
#         logger.info("Added %d STS+mask pages", added_pages)
#     else:
#         logger.warning("sts or mask directory not found, skipping STS pages.")

#     ppt_path = output_folder / (output_folder.name + ".pptx")
#     prs.save(str(ppt_path))
#     logger.info("PPT saved to %s", ppt_path)


# def generate_ppt(output_folder: Path, prs: Presentation) -> None:
#     """为单个输出文件夹生成幻灯片，追加到已有的 prs 中。"""
#     # ---------- 首页：使用 output_folder 的父级三层路径和创建时间 ----------
#     slide = prs.slides.add_slide(prs.slide_layouts[6])

#     parts = output_folder.resolve().parts
#     name_folder = "\\".join(parts[-3:]) if len(parts) >= 3 else "\\".join(parts)
#     time_folder = str(datetime.fromtimestamp(output_folder.resolve().stat().st_ctime))
#     time_folder = time_folder.split(".")[0]

#     text_box1 = slide.shapes.add_textbox(Cm(4), Cm(5.5), Inches(4), Inches(1))
#     tf1 = text_box1.text_frame
#     p1 = tf1.add_paragraph()
#     p1.text = name_folder
#     p1.font.size = Pt(44)

#     line = slide.shapes.add_shape(1, Cm(0), Cm(9.5), Cm(25.4), Cm(1.3))
#     line.fill.solid()
#     line.fill.fore_color.rgb = RGBColor(0, 0, 255)
#     line.line.fill.background()

#     text_box2 = slide.shapes.add_textbox(Cm(4), Cm(13), Inches(4), Inches(1))
#     tf2 = text_box2.text_frame
#     p2 = tf2.add_paragraph()
#     p2.text = time_folder
#     p2.font.size = Pt(36)

#     # ---------- 内容页：Z_mtrx图片 ----------
#     img_dir = output_folder / "images"
#     if not img_dir.is_dir():
#         logger.error("Images directory not found: %s", img_dir)
#         return

#     z_images = get_z_images_sorted(img_dir)
#     if z_images:
#         pic_width = Cm(4.5)
#         pic_height = Cm(4.5)
#         left_margin = Cm(0.5)
#         gap_x = (Cm(25.4) - 2 * left_margin - 5 * pic_width) / 4
#         row1_y = Cm(2)
#         row2_y = Cm(9)
#         label_height = Cm(0.5)
#         label_offset = Cm(0.6)

#         page_count = 0
#         for start in range(0, len(z_images), 10):
#             chunk = z_images[start : start + 10]
#             slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用传入的 prs
#             for i, img_path in enumerate(chunk):
#                 col = i % 5
#                 row = i // 5
#                 left = left_margin + col * (pic_width + gap_x)
#                 top = row1_y if row == 0 else row2_y

#                 # 读取原始尺寸，计算缩放以适应 4.5x4.5 cm² 区域，保持比例
#                 from PIL import Image

#                 with Image.open(str(img_path)) as im:
#                     w_px, h_px = im.size
#                 scale = min(4.5 / w_px, 4.5 / h_px)
#                 display_w = Cm(w_px * scale)
#                 display_h = Cm(h_px * scale)

#                 slide.shapes.add_picture(
#                     str(img_path), left, top, width=display_w, height=display_h
#                 )

#                 # 提取短标签
#                 stem = img_path.stem
#                 label = stem.split("--")[-1].replace(".Z_mtrx_forward_up", "")

#                 # 图片正上方添加标签文本框
#                 txBox = slide.shapes.add_textbox(
#                     left, top - label_offset, pic_width, label_height
#                 )
#                 tf = txBox.text_frame
#                 tf.word_wrap = True
#                 p = tf.paragraphs[0]
#                 p.text = label
#                 p.font.size = Pt(16)
#                 p.alignment = PP_ALIGN.CENTER

#                 # 图片下方参数信息
#                 if stem.endswith(".Z_mtrx_forward_up"):
#                     orig_name = stem[: -len(".Z_mtrx_forward_up")]
#                     h5_path = output_folder / (orig_name + ".Z_mtrx.h5")
#                     text_str = ""
#                     if h5_path.exists():
#                         try:
#                             with h5py.File(h5_path, "r") as hf:
#                                 params = hf["parameters"]
#                                 bias = round(
#                                     float(
#                                         params["GapVoltageControl.Voltage"].attrs[
#                                             "value"
#                                         ]
#                                     ),
#                                     3,
#                                 )
#                                 sp_pa = round(
#                                     float(params["Regulator.Setpoint_1"].attrs["value"])
#                                     * 1e12,
#                                     2,
#                                 )
#                                 img_grp = hf["images"]
#                                 first_dir = next(iter(img_grp))
#                                 w_m = float(img_grp[first_dir].attrs["width"])
#                                 h_m = float(img_grp[first_dir].attrs["height"])
#                                 real_w_nm = w_m * 1e9
#                                 real_h_nm = h_m * 1e9
#                                 text_str = (
#                                     f"{bias} V, {sp_pa} pA\n"
#                                     f"{real_w_nm:.1f} × {real_h_nm:.1f} nm²"
#                                 )
#                         except Exception as e:
#                             logger.warning(
#                                 "Failed to read params from %s: %s", h5_path, e
#                             )
#                     else:
#                         logger.warning("HDF5 not found: %s", h5_path)
#                     if text_str:
#                         txBox_bottom = slide.shapes.add_textbox(
#                             left,
#                             top + pic_height + Cm(0.05),
#                             pic_width,
#                             Cm(0.6),
#                         )
#                         tf_bottom = txBox_bottom.text_frame
#                         tf_bottom.word_wrap = True
#                         p_bottom = tf_bottom.paragraphs[0]
#                         p_bottom.text = text_str
#                         p_bottom.font.size = Pt(14)
#                         p_bottom.alignment = PP_ALIGN.CENTER

#             page_count += 1
#         logger.info("Added %d content pages with Z_mtrx images", page_count)
#     else:
#         logger.warning("No Z_mtrx_forward_up images found in %s", img_dir)

#     # ---------- STS 与 mask 叠加页 ----------
#     sts_dir = output_folder / "sts"
#     mask_dir = output_folder / "mask"
#     if sts_dir.is_dir() and mask_dir.is_dir():
#         sts_files = sorted(sts_dir.glob("*.png"), key=lambda p: natural_key(p.stem))
#         added_pages = 0
#         for sts_path in sts_files:
#             info = parse_sts_filename(sts_path.name)
#             if info is None:
#                 continue
#             img_name, group_idx = info
#             img_path = img_dir / f"{img_name}_forward_up.png"
#             mask_path = mask_dir / f"{img_name}_group{group_idx}.png"
#             if not img_path.exists() or not mask_path.exists():
#                 logger.warning("Missing image/mask for %s", sts_path.name)
#                 continue

#             slide = prs.slides.add_slide(prs.slide_layouts[6])

#             left_img = Cm(1)
#             top_img = Cm(3)
#             img_size = Cm(12)

#             slide.shapes.add_picture(
#                 str(img_path), left_img, top_img, width=img_size, height=img_size
#             )
#             slide.shapes.add_picture(
#                 str(mask_path), left_img, top_img, width=img_size, height=img_size
#             )

#             sts_left = Cm(14)
#             sts_top = Cm(3)
#             sts_width = Cm(9)
#             sts_height = Cm(9)
#             slide.shapes.add_picture(
#                 str(sts_path), sts_left, sts_top, width=sts_width, height=sts_height
#             )

#             txBox = slide.shapes.add_textbox(left_img, Cm(1.5), img_size, Cm(1))
#             tf = txBox.text_frame
#             p = tf.paragraphs[0]
#             p.text = f"{img_name}  Group {group_idx}"
#             p.font.size = Pt(24)
#             p.alignment = PP_ALIGN.LEFT

#             added_pages += 1
#         logger.info("Added %d STS+mask pages", added_pages)
#     else:
#         logger.warning("sts or mask directory not found, skipping STS pages.")


def generate_ppt(output_folder: Path, prs: Presentation) -> None:
    """为单个输出文件夹生成幻灯片，追加到已有的 prs 中。"""
    # ---------- 首页：使用 output_folder 的父级三层路径和创建时间 ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    parts = output_folder.resolve().parts
    name_folder = "\\".join(parts[-3:]) if len(parts) >= 3 else "\\".join(parts)
    time_folder = str(datetime.fromtimestamp(output_folder.resolve().stat().st_ctime))
    time_folder = time_folder.split(".")[0]

    text_box1 = slide.shapes.add_textbox(Cm(4), Cm(5.5), Inches(4), Inches(1))
    tf1 = text_box1.text_frame
    p1 = tf1.add_paragraph()
    p1.text = name_folder
    p1.font.size = Pt(44)

    line = slide.shapes.add_shape(1, Cm(0), Cm(9.5), Cm(25.4), Cm(1.3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 255)
    line.line.fill.background()

    text_box2 = slide.shapes.add_textbox(Cm(4), Cm(13), Inches(4), Inches(1))
    tf2 = text_box2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = time_folder
    p2.font.size = Pt(36)

    # ---------- 内容页：Z_mtrx图片 ----------
    img_dir = output_folder / "images"
    if not img_dir.is_dir():
        logger.error("Images directory not found: %s", img_dir)
        return

    z_images = get_z_images_sorted(img_dir)
    if z_images:
        pic_width = Cm(4.5)
        pic_height = Cm(4.5)
        left_margin = Cm(0.5)
        gap_x = (Cm(25.4) - 2 * left_margin - 5 * pic_width) / 4
        row1_y = Cm(2)
        row2_y = Cm(9)
        label_height = Cm(0.5)
        label_offset = Cm(0.6)

        page_count = 0
        for start in range(0, len(z_images), 10):
            chunk = z_images[start : start + 10]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            for i, img_path in enumerate(chunk):
                col = i % 5
                row = i // 5
                left = left_margin + col * (pic_width + gap_x)
                top = row1_y if row == 0 else row2_y

                stem = img_path.stem

                # ----- 提前读取 H5 参数，获取角度等信息 -----
                text_str = ""
                angle = 0.0
                if stem.endswith(".Z_mtrx_forward_up"):
                    orig_name = stem[: -len(".Z_mtrx_forward_up")]
                    h5_path = output_folder / (orig_name + ".Z_mtrx.h5")
                    if h5_path.exists():
                        try:
                            with h5py.File(h5_path, "r") as hf:
                                params = hf["parameters"]
                                bias = round(
                                    float(
                                        params["GapVoltageControl.Voltage"].attrs[
                                            "value"
                                        ]
                                    ),
                                    3,
                                )
                                sp_pa = round(
                                    float(params["Regulator.Setpoint_1"].attrs["value"])
                                    * 1e12,
                                    2,
                                )
                                img_grp = hf["images"]
                                first_dir = next(iter(img_grp))
                                w_m = float(img_grp[first_dir].attrs["width"])
                                h_m = float(img_grp[first_dir].attrs["height"])
                                angle = float(img_grp[first_dir].attrs["angle"])
                                real_w_nm = w_m * 1e9
                                real_h_nm = h_m * 1e9
                                text_str = (
                                    f"{bias} V, {sp_pa} pA\n"
                                    f"{real_w_nm:.1f} × {real_h_nm:.1f} nm\u00b2 ({angle:.1f}\u00b0)"
                                )
                        except Exception as e:
                            logger.warning(
                                "Failed to read params from %s: %s", h5_path, e
                            )
                    else:
                        logger.warning("HDF5 not found: %s", h5_path)

                # ----- 插入图片（保持比例、居中、可旋转）-----
                with Image.open(str(img_path)) as im:
                    w_px, h_px = im.size
                scale = min(4.5 / w_px, 4.5 / h_px)
                display_w = Cm(w_px * scale)
                display_h = Cm(h_px * scale)

                # 计算居中偏移量
                offset_x = (pic_width - display_w) / 2
                offset_y = (pic_height - display_h) / 2
                picture = slide.shapes.add_picture(
                    str(img_path),
                    left + offset_x,
                    top + offset_y,
                    width=display_w,
                    height=display_h,
                )
                picture.rotation = angle

                # ----- 标签与底部参数文本 -----
                label = stem.split("--")[-1].replace(".Z_mtrx_forward_up", "")

                txBox = slide.shapes.add_textbox(
                    left, top - label_offset, pic_width, label_height
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.CENTER

                if text_str:
                    txBox_bottom = slide.shapes.add_textbox(
                        left,
                        top + pic_height + Cm(0.05),
                        pic_width,
                        Cm(0.6),
                    )
                    tf_bottom = txBox_bottom.text_frame
                    tf_bottom.word_wrap = True
                    p_bottom = tf_bottom.paragraphs[0]
                    p_bottom.text = text_str
                    p_bottom.font.size = Pt(14)
                    p_bottom.alignment = PP_ALIGN.CENTER

            page_count += 1
        logger.info("Added %d content pages with Z_mtrx images", page_count)
    else:
        logger.warning("No Z_mtrx_forward_up images found in %s", img_dir)

    # ---------- STS 与 mask 叠加页 ----------
    sts_dir = output_folder / "sts"
    mask_dir = output_folder / "mask"
    if sts_dir.is_dir() and mask_dir.is_dir():
        sts_files = sorted(sts_dir.glob("*.png"), key=lambda p: natural_key(p.stem))
        added_pages = 0
        for sts_path in sts_files:
            info = parse_sts_filename(sts_path.name)
            if info is None:
                continue
            img_name, group_idx = info
            img_path = img_dir / f"{img_name}_forward_up.png"
            mask_path = mask_dir / f"{img_name}_group{group_idx}.png"
            if not img_path.exists() or not mask_path.exists():
                logger.warning("Missing image/mask for %s", sts_path.name)
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            left_img = Cm(1)
            top_img = Cm(3)
            img_size = Cm(12)

            slide.shapes.add_picture(
                str(img_path), left_img, top_img, width=img_size, height=img_size
            )
            slide.shapes.add_picture(
                str(mask_path), left_img, top_img, width=img_size, height=img_size
            )

            sts_left = Cm(14)
            sts_top = Cm(3)
            sts_width = Cm(9)
            sts_height = Cm(9)
            slide.shapes.add_picture(
                str(sts_path), sts_left, sts_top, width=sts_width, height=sts_height
            )

            txBox = slide.shapes.add_textbox(left_img, Cm(1.5), img_size, Cm(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{img_name}  Group {group_idx}"
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.LEFT

            added_pages += 1
        logger.info("Added %d STS+mask pages", added_pages)
    else:
        logger.warning("sts or mask directory not found, skipping STS pages.")


# def main() -> None:
#     logging.basicConfig(
#         level=logging.INFO,
#         format="%(asctime)s [%(levelname)s] %(message)s",
#         datefmt="%Y-%m-%d %H:%M:%S",
#         stream=sys.stdout,
#     )

#     if len(sys.argv) < 2:
#         print("Usage: python generate_ppt.py <output_folder>")
#         sys.exit(1)

#     output_folder = Path(sys.argv[1])
#     if not output_folder.is_dir():
#         logger.error("'%s' is not a directory", output_folder)
#         sys.exit(1)

#     generate_ppt(output_folder)


def main() -> None:
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s [%(levelname)s] %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
        stream=sys.stdout,
    )

    if len(sys.argv) < 2:
        print("Usage: python preapre_ppt.py <parent_folder>")
        sys.exit(1)

    parent = Path(sys.argv[1])
    if not parent.is_dir():
        logger.error("'%s' is not a directory", parent)
        sys.exit(1)

    subdirs = sorted(
        [d for d in parent.iterdir() if d.is_dir()],
        key=lambda d: natural_key(d.name),
    )

    if not subdirs:
        logger.error("No subdirectories found in %s", parent)
        sys.exit(1)

    prs = Presentation()
    for sub in subdirs:
        try:
            generate_ppt(sub, prs)
        except Exception as e:
            logger.error("Failed to process %s: %s", sub, e)

    ppt_path = parent / "combined.pptx"
    prs.save(str(ppt_path))
    logger.info("Combined PPT saved to %s", ppt_path)


if __name__ == "__main__":
    main()
